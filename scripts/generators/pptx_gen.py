"""
Clase base para generar presentaciones (.pptx) y scripts de teleprompter (.md).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from scripts.config.constants import (
    ColorRGB, CURSO_NOMBRE, INSTRUCTOR, INSTRUCTOR_TITULO, ANIO
)

# ── Colores pptx ─────────────────────────────────────────────────
PPTX_AZUL = RGBColor(*ColorRGB.AZUL)
PPTX_VERDE = RGBColor(*ColorRGB.VERDE)
PPTX_ROJO = RGBColor(*ColorRGB.ROJO)
PPTX_AMARILLO = RGBColor(*ColorRGB.AMARILLO)
PPTX_TEXTO = RGBColor(*ColorRGB.TEXTO_OSCURO)
PPTX_TEXTO_MEDIO = RGBColor(*ColorRGB.TEXTO_MEDIO)
PPTX_BLANCO = RGBColor(*ColorRGB.BLANCO)
PPTX_FONDO = RGBColor(*ColorRGB.FONDO_CLARO)


class SlidesGenerator:
    """Base class for generating PPTX presentations with teleprompter script."""

    def __init__(self, filename: str, output_dir: Path, script_filename: str = "",
                 script_dir: Path = None):
        self.filename = filename
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.script_lines = []
        self.script_filename = script_filename
        self.script_dir = script_dir or output_dir
        self.script_dir.mkdir(parents=True, exist_ok=True)
        self.slide_count = 0

    @property
    def filepath(self) -> Path:
        return self.output_dir / self.filename

    @property
    def script_filepath(self) -> Path:
        return self.script_dir / self.script_filename

    def _add_bg(self, slide, color_rgb=ColorRGB.FONDO_CLARO):
        """Set solid background color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color_rgb)

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT,
                     font_name="Calibri"):
        """Add a text box to a slide."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or PPTX_TEXTO
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_bullet_list(self, slide, left, top, width, height, items,
                         font_size=16, color=None):
        """Add a bulleted list to a slide."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color or PPTX_TEXTO
            p.font.name = "Calibri"
            p.space_after = Pt(6)
        return txBox

    def _add_shape_rect(self, slide, left, top, width, height, fill_rgb, text="",
                        font_size=14, font_color=None, bold=False):
        """Add a filled rectangle with optional text."""
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
        shape.line.fill.background()
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = font_color or PPTX_BLANCO
            p.font.name = "Calibri"
        return shape

    def add_title_slide(self, modulo_num, modulo_nombre, subtitulo=""):
        """Add a module title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        self._add_bg(slide, ColorRGB.AZUL)
        self.slide_count += 1

        # Module number
        self._add_textbox(slide, 0.5, 1.0, 12, 1.0,
                          f"MÓDULO {modulo_num}" if modulo_num else "",
                          font_size=20, bold=True, color=PPTX_BLANCO,
                          alignment=PP_ALIGN.CENTER)
        # Title
        self._add_textbox(slide, 0.5, 2.2, 12, 2.0,
                          modulo_nombre, font_size=36, bold=True,
                          color=PPTX_BLANCO, alignment=PP_ALIGN.CENTER)
        # Subtitle
        if subtitulo:
            self._add_textbox(slide, 1, 4.5, 11, 1.0,
                              subtitulo, font_size=18,
                              color=PPTX_BLANCO, alignment=PP_ALIGN.CENTER)
        # Footer
        self._add_textbox(slide, 0.5, 6.2, 12, 0.5,
                          f"{INSTRUCTOR} — {INSTRUCTOR_TITULO} — {CURSO_NOMBRE} {ANIO}",
                          font_size=12, color=PPTX_BLANCO,
                          alignment=PP_ALIGN.CENTER)

        self.script_lines.append(f"# Módulo {modulo_num}: {modulo_nombre}\n")
        self.script_lines.append(f"## Slide 1 — Portada\n")
        return slide

    def add_content_slide(self, title, bullets=None, body_text=None,
                          script_text="", note=""):
        """Add a standard content slide with title + bullets or body text."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self.slide_count += 1

        # Blue accent bar
        self._add_shape_rect(slide, 0, 0, 0.15, 7.5, ColorRGB.AZUL)

        # Title
        self._add_textbox(slide, 0.6, 0.3, 12, 0.8,
                          title, font_size=28, bold=True, color=PPTX_AZUL)

        # Content
        if bullets:
            self._add_bullet_list(slide, 0.8, 1.5, 11.5, 5.0, bullets, font_size=18)
        elif body_text:
            self._add_textbox(slide, 0.8, 1.5, 11.5, 5.0,
                              body_text, font_size=18)

        # Script
        if script_text:
            self.script_lines.append(f"## Slide {self.slide_count} — {title}\n")
            self.script_lines.append(f"{script_text}\n")

        return slide

    def add_closing_slide(self, next_module="", resources=None):
        """Add a closing/resources slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, ColorRGB.AZUL)
        self.slide_count += 1

        self._add_textbox(slide, 0.5, 1.5, 12, 1.5,
                          "Recursos y Siguiente Paso",
                          font_size=32, bold=True, color=PPTX_BLANCO,
                          alignment=PP_ALIGN.CENTER)

        items = resources or []
        if next_module:
            items.append(f"Siguiente: {next_module}")

        if items:
            self._add_bullet_list(slide, 1.5, 3.5, 10, 3.0, items,
                                  font_size=18, color=PPTX_BLANCO)

        self._add_textbox(slide, 0.5, 6.5, 12, 0.5,
                          f"{CURSO_NOMBRE} — {INSTRUCTOR}",
                          font_size=12, color=PPTX_BLANCO,
                          alignment=PP_ALIGN.CENTER)

        self.script_lines.append(f"## Slide {self.slide_count} — Cierre\n")
        return slide

    def save(self):
        """Save PPTX and script."""
        self.prs.save(str(self.filepath))
        print(f"  ✓ {self.filepath.name}")

        if self.script_filename and self.script_lines:
            script_path = self.script_filepath
            script_path.write_text("\n".join(self.script_lines), encoding="utf-8")
            print(f"  ✓ {script_path.name}")

        return self.filepath
